from __future__ import annotations

from dataclasses import dataclass
from html import escape
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
SCREENSHOT_DIR = ROOT / "screenshots" / "compact_demo"

DRAWIO_PATH = DOCS_DIR / "inkpi-project-flow.drawio"
PNG_PATH = DOCS_DIR / "inkpi-project-flow.png"
PPTX_PATH = DOCS_DIR / "inkpi-project-overview.pptx"

CANVAS_W = 2200
CANVAS_H = 1400

BG = "#F6F1E8"
PANEL_BG = "#FFF8EF"
TITLE_BG = "#1C1816"
TEXT = "#241E1A"
MUTED = "#6B5F55"
LINE = "#8A8076"


@dataclass(frozen=True)
class Section:
    id: str
    label: str
    x: int
    y: int
    w: int
    h: int
    fill: str
    stroke: str


@dataclass(frozen=True)
class Node:
    id: str
    title: str
    body: str
    x: int
    y: int
    w: int
    h: int
    fill: str
    stroke: str
    font_color: str = TEXT


@dataclass(frozen=True)
class Edge:
    id: str
    source: str
    target: str
    label: str = ""
    dashed: bool = False


SECTIONS = [
    Section("sec_input", "1. 输入与入口", 40, 110, 260, 42, "#E6EDF8", "#8AA0BE"),
    Section("sec_runtime", "2. 本地评测主链", 340, 110, 1340, 42, "#FCE9D6", "#D79B4B"),
    Section("sec_cloud", "3. 存储 / 同步 / 云端", 1710, 110, 450, 42, "#F8D9D7", "#B55B54"),
]


NODES = [
    Node(
        "desktop",
        "PyQt 正式界面",
        "main.py\nviews/home_view.py\nviews/camera_view.py\nviews/result_view.py\n主分 + 四维解释分，不显示 debug",
        50,
        190,
        220,
        130,
        "#D9E8FB",
        "#6C8EBF",
    ),
    Node(
        "webui",
        "WebUI 调试端",
        "web_ui/app.py\nindex.html + app.js\n拍照 / 上传 / 历史 / 详情\n详情页显示 score_debug",
        50,
        350,
        220,
        140,
        "#D9E8FB",
        "#6C8EBF",
    ),
    Node(
        "input",
        "图像输入",
        "摄像头预览 / 拍照\n本地图片上传\n统一进入单字评测链路",
        50,
        520,
        220,
        100,
        "#FFF0C9",
        "#D3AA4A",
    ),
    Node(
        "preprocess",
        "PreprocessingService",
        "预检、透视校正、去网格、主体提取\n产出 processed_image + ocr_image",
        360,
        290,
        270,
        130,
        "#FFE6CC",
        "#D79B00",
    ),
    Node(
        "ocr",
        "LocalOcrService",
        "PaddleOCR 本地识别\n输出 character + confidence",
        710,
        170,
        250,
        110,
        "#DDEFD9",
        "#82B366",
    ),
    Node(
        "quality",
        "QualityScorerService",
        "ONNX 主分模型\n输出 total_score / quality_level / quality_confidence\n并保留 probabilities / quality_features / calibration",
        1040,
        150,
        330,
        150,
        "#DDEFD9",
        "#82B366",
    ),
    Node(
        "dimension",
        "DimensionScorerService",
        "复用 quality_features + geometry_features + calibration\n输出 structure / stroke / integrity / stability\n并生成 strongest / weakest 解释",
        1040,
        360,
        330,
        160,
        "#EEF6EA",
        "#82B366",
    ),
    Node(
        "evaluation",
        "EvaluationService",
        "编排 processed_image + ocr_image\n调用 OCR / 主分 / 四维解释层\n生成统一 EvaluationResult",
        710,
        330,
        260,
        130,
        "#E7D7EE",
        "#9673A6",
    ),
    Node(
        "result",
        "EvaluationResult",
        "total_score / quality_level / feedback\ncharacter_name / quality_confidence\ndimension_scores / score_debug",
        710,
        550,
        300,
        130,
        "#E7D7EE",
        "#9673A6",
    ),
    Node(
        "database",
        "DatabaseService",
        "本地 SQLite\n新增 dimension_scores_json\n新增 score_debug_json\n旧记录保持 null",
        1080,
        560,
        270,
        140,
        "#F8D4D3",
        "#B85450",
    ),
    Node(
        "sync",
        "CloudSyncService",
        "upload_result_async()\n上传主分、四维分和 debug JSON\nPOST /api/device/results",
        1420,
        570,
        250,
        120,
        "#F8D4D3",
        "#B85450",
    ),
    Node(
        "qt_result",
        "Qt 结果页",
        "正式展示端\n显示主分、等级、反馈、四维卡片\n不显示 score_debug",
        350,
        860,
        260,
        130,
        "#D9E8FB",
        "#6C8EBF",
    ),
    Node(
        "webui_debug",
        "WebUI 详情页",
        "调试展示端\n历史列表显示四维摘要\n详情页展开 probabilities / quality_features /\ngeometry_features / calibration",
        680,
        850,
        360,
        160,
        "#D9E8FB",
        "#6C8EBF",
    ),
    Node(
        "miniapp",
        "MiniApp 结果详情",
        "小程序只显示四维解释分和主反馈\n历史列表保持轻量\n不展示 score_debug",
        1100,
        860,
        300,
        130,
        "#DDEFD9",
        "#82B366",
    ),
    Node(
        "cloud_api",
        "Cloud API",
        "cloud_api/app.py\n列表接口返回 dimension_scores\n详情接口返回 dimension_scores + score_debug",
        1780,
        250,
        320,
        150,
        "#E7D7EE",
        "#9673A6",
    ),
    Node(
        "cloud_db",
        "CloudDatabase",
        "cloud_api/storage.py\nresults 表新增\n dimension_scores_json / score_debug_json",
        1780,
        450,
        320,
        130,
        "#E7D7EE",
        "#9673A6",
    ),
    Node(
        "training",
        "训练链路",
        "build_quality_manifest.py\ntrain_quality_scorer.py\n导出 quality_scorer.onnx\n持续迭代主分模型",
        1780,
        780,
        320,
        130,
        "#FFF0C9",
        "#D3AA4A",
    ),
    Node(
        "model_artifact",
        "模型与静态资产",
        "models/quality_scorer.onnx\nmodels/quality_scorer.metrics.json\ndocs/inkpi-project-flow.png",
        1780,
        950,
        320,
        110,
        "#FFF0C9",
        "#D3AA4A",
    ),
    Node(
        "main_score_note",
        "主分校准细节",
        "quality_features:\nfg_ratio / bbox_ratio / center_quality /\ncomponent_norm / edge_touch / texture_std\n\ncalibration:\nfeature_quality / probability_margin /\nquality_confidence_norm / score_range_fit",
        360,
        500,
        270,
        220,
        "#F7F3E6",
        "#B9AC8F",
    ),
    Node(
        "dimension_note",
        "四维解释分规则",
        "结构 = center_quality + bbox_ratio_band + projection_balance + bbox_fill_band\n笔画 = texture_std_band + orientation_concentration + ink_ratio_band + component_norm_band\n完整 = ocr_confidence + dominant_share + (1-edge_touch) + subject_edge_safe\n稳定 = quality_confidence + probability_margin + feature_quality + score_range_fit",
        360,
        760,
        1180,
        150,
        "#F7F3E6",
        "#B9AC8F",
    ),
]


EDGES = [
    Edge("e1", "desktop", "preprocess"),
    Edge("e2", "webui", "preprocess"),
    Edge("e3", "input", "preprocess"),
    Edge("e4", "preprocess", "ocr", "ocr_image"),
    Edge("e5", "preprocess", "quality", "processed_image"),
    Edge("e6", "preprocess", "dimension", "geometry ROI"),
    Edge("e7", "ocr", "evaluation", "character + confidence"),
    Edge("e8", "quality", "evaluation", "主分 + 校准"),
    Edge("e9", "dimension", "evaluation", "四维解释分"),
    Edge("e10", "evaluation", "result"),
    Edge("e11", "result", "database"),
    Edge("e12", "database", "sync"),
    Edge("e13", "sync", "cloud_api", "/api/device/results"),
    Edge("e14", "cloud_api", "cloud_db"),
    Edge("e15", "database", "qt_result"),
    Edge("e16", "database", "webui_debug"),
    Edge("e17", "cloud_api", "miniapp"),
    Edge("e18", "training", "model_artifact"),
    Edge("e19", "model_artifact", "quality", "quality_scorer.onnx", True),
    Edge("e20", "main_score_note", "quality"),
    Edge("e21", "dimension_note", "dimension"),
]


def drawio_color_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def wrap_text(text: str, font: ImageFont.FreeTypeFont, width: int) -> list[str]:
    draw = ImageDraw.Draw(Image.new("RGB", (10, 10)))
    lines: list[str] = []
    for paragraph in text.split("\n"):
        words = paragraph.split(" ")
        if len(words) == 1 and any("\u4e00" <= ch <= "\u9fff" for ch in paragraph):
            current = ""
            for ch in paragraph:
                candidate = f"{current}{ch}"
                if draw.textbbox((0, 0), candidate, font=font)[2] <= width or not current:
                    current = candidate
                else:
                    lines.append(current)
                    current = ch
            lines.append(current)
            continue

        current = ""
        for word in words:
            candidate = word if not current else f"{current} {word}"
            if draw.textbbox((0, 0), candidate, font=font)[2] <= width or not current:
                current = candidate
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def get_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        Path(r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def build_drawio() -> None:
    mxfile = ET.Element(
        "mxfile",
        host="app.diagrams.net",
        agent="Codex",
        version="29.6.6",
        type="device",
    )
    diagram = ET.SubElement(mxfile, "diagram", id="inkpi-project-flow", name="InkPi Project Flow")
    graph = ET.SubElement(
        diagram,
        "mxGraphModel",
        dx="2377",
        dy="1238",
        grid="1",
        gridSize="10",
        guides="1",
        tooltips="1",
        connect="1",
        arrows="1",
        fold="1",
        page="1",
        pageScale="1",
        pageWidth="2200",
        pageHeight="1400",
        math="0",
        shadow="0",
    )
    root = ET.SubElement(graph, "root")
    ET.SubElement(root, "mxCell", id="0")
    ET.SubElement(root, "mxCell", id="1", parent="0")

    title = ET.SubElement(
        root,
        "mxCell",
        id="title",
        parent="1",
        value=(
            "<b>InkPi 当前项目全流程（含四维解释分）</b><br>"
            "主链：入口 → 预处理 → OCR / 主分评分 / 四维解释 → EvaluationResult → 本地落库 → 云端同步"
        ),
        style=(
            "rounded=1;whiteSpace=wrap;html=1;fillColor=#1C1816;strokeColor=#1C1816;"
            "fontColor=#FFFFFF;fontStyle=1;fontSize=22;align=center;"
        ),
        vertex="1",
    )
    ET.SubElement(title, "mxGeometry", {"x": "20", "y": "20", "width": "2160", "height": "70", "as": "geometry"})

    for section in SECTIONS:
        cell = ET.SubElement(
            root,
            "mxCell",
            id=section.id,
            parent="1",
            value=escape(section.label),
            style=(
                f"rounded=1;whiteSpace=wrap;html=1;fillColor={section.fill};strokeColor={section.stroke};"
                "fontStyle=1;fontSize=18;align=center;"
            ),
            vertex="1",
        )
        ET.SubElement(
            cell,
            "mxGeometry",
            {
                "x": str(section.x),
                "y": str(section.y),
                "width": str(section.w),
                "height": str(section.h),
                "as": "geometry",
            },
        )

    for node in NODES:
        value = f"<b>{escape(node.title)}</b><br>{escape(node.body).replace(chr(10), '<br>')}"
        cell = ET.SubElement(
            root,
            "mxCell",
            id=node.id,
            parent="1",
            value=value,
            style=(
                f"rounded=1;whiteSpace=wrap;html=1;fillColor={node.fill};strokeColor={node.stroke};"
                f"fontSize=14;fontColor={node.font_color};"
            ),
            vertex="1",
        )
        ET.SubElement(
            cell,
            "mxGeometry",
            {
                "x": str(node.x),
                "y": str(node.y),
                "width": str(node.w),
                "height": str(node.h),
                "as": "geometry",
            },
        )

    for edge in EDGES:
        cell = ET.SubElement(
            root,
            "mxCell",
            id=edge.id,
            parent="1",
            source=edge.source,
            target=edge.target,
            value=escape(edge.label),
            style=(
                "edgeStyle=orthogonalEdgeStyle;rounded=0;orthogonalLoop=1;jettySize=auto;html=1;"
                f"strokeColor=#666666;strokeWidth=2;endArrow=block;endFill=1;"
                f"{'dashed=1;' if edge.dashed else ''}"
            ),
            edge="1",
        )
        ET.SubElement(cell, "mxGeometry", {"relative": "1", "as": "geometry"})

    ET.indent(mxfile)
    DRAWIO_PATH.write_text(ET.tostring(mxfile, encoding="unicode"), encoding="utf-8")


def center_left(node: Node) -> tuple[int, int]:
    return node.x, node.y + node.h // 2


def center_right(node: Node) -> tuple[int, int]:
    return node.x + node.w, node.y + node.h // 2


def center_top(node: Node) -> tuple[int, int]:
    return node.x + node.w // 2, node.y


def center_bottom(node: Node) -> tuple[int, int]:
    return node.x + node.w // 2, node.y + node.h


def resolve_points(source: Node, target: Node) -> tuple[tuple[int, int], tuple[int, int]]:
    if source.x + source.w < target.x:
        return center_right(source), center_left(target)
    if target.x + target.w < source.x:
        return center_left(source), center_right(target)
    if source.y + source.h < target.y:
        return center_bottom(source), center_top(target)
    return center_top(source), center_bottom(target)


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str, dashed: bool = False) -> None:
    points = [start]
    if abs(end[0] - start[0]) > abs(end[1] - start[1]):
        mid_x = (start[0] + end[0]) // 2
        points.extend([(mid_x, start[1]), (mid_x, end[1])])
    else:
        mid_y = (start[1] + end[1]) // 2
        points.extend([(start[0], mid_y), (end[0], mid_y)])
    points.append(end)

    if dashed:
        for left, right in zip(points, points[1:]):
            draw_dashed_line(draw, left, right, color)
    else:
        draw.line(points, fill=color, width=4, joint="curve")

    arrow_size = 12
    ex, ey = end
    if len(points) >= 2:
        px, py = points[-2]
    else:
        px, py = start
    if abs(ex - px) >= abs(ey - py):
        if ex >= px:
            polygon = [(ex, ey), (ex - arrow_size, ey - 6), (ex - arrow_size, ey + 6)]
        else:
            polygon = [(ex, ey), (ex + arrow_size, ey - 6), (ex + arrow_size, ey + 6)]
    else:
        if ey >= py:
            polygon = [(ex, ey), (ex - 6, ey - arrow_size), (ex + 6, ey - arrow_size)]
        else:
            polygon = [(ex, ey), (ex - 6, ey + arrow_size), (ex + 6, ey + arrow_size)]
    draw.polygon(polygon, fill=color)


def draw_dashed_line(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str) -> None:
    dash = 12
    gap = 8
    if start[0] == end[0]:
        y0, y1 = sorted([start[1], end[1]])
        x = start[0]
        current = y0
        while current < y1:
            draw.line([(x, current), (x, min(current + dash, y1))], fill=color, width=4)
            current += dash + gap
    else:
        x0, x1 = sorted([start[0], end[0]])
        y = start[1]
        current = x0
        while current < x1:
            draw.line([(current, y), (min(current + dash, x1), y)], fill=color, width=4)
            current += dash + gap


def draw_multiline_text(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    width: int,
    title: str,
    body: str,
    title_font: ImageFont.FreeTypeFont,
    body_font: ImageFont.FreeTypeFont,
    color: str,
) -> None:
    title_box = draw.textbbox((0, 0), title, font=title_font)
    draw.text((x, y), title, font=title_font, fill=color)
    current_y = y + (title_box[3] - title_box[1]) + 10
    for line in body.split("\n"):
        wrapped = wrap_text(line, body_font, width)
        for wrapped_line in wrapped:
            draw.text((x, current_y), wrapped_line, font=body_font, fill=color)
            line_box = draw.textbbox((0, 0), wrapped_line, font=body_font)
            current_y += (line_box[3] - line_box[1]) + 4


def build_png() -> None:
    image = Image.new("RGB", (CANVAS_W, CANVAS_H), BG)
    draw = ImageDraw.Draw(image)

    title_font = get_font(34, bold=True)
    section_font = get_font(24, bold=True)
    node_title_font = get_font(26, bold=True)
    body_font = get_font(20)
    small_font = get_font(18)

    draw.rounded_rectangle((20, 20, CANVAS_W - 20, 90), radius=22, fill=TITLE_BG, outline=TITLE_BG)
    draw.text((60, 36), "InkPi 当前项目全流程（含四维解释分）", font=title_font, fill="#FFFFFF")
    draw.text(
        (60, 74),
        "主链：入口 → 预处理 → OCR / 主分评分 / 四维解释 → EvaluationResult → 本地落库 → 云端同步",
        font=small_font,
        fill="#DDD4CA",
    )

    for section in SECTIONS:
        draw.rounded_rectangle(
            (section.x, section.y, section.x + section.w, section.y + section.h),
            radius=18,
            fill=section.fill,
            outline=section.stroke,
            width=3,
        )
        box = draw.textbbox((0, 0), section.label, font=section_font)
        draw.text(
            (section.x + (section.w - (box[2] - box[0])) / 2, section.y + 6),
            section.label,
            font=section_font,
            fill=TEXT,
        )

    node_map = {node.id: node for node in NODES}

    for edge in EDGES:
        start, end = resolve_points(node_map[edge.source], node_map[edge.target])
        draw_arrow(draw, start, end, LINE, dashed=edge.dashed)
        if edge.label:
            label_x = (start[0] + end[0]) // 2
            label_y = (start[1] + end[1]) // 2 - 16
            label_box = draw.textbbox((0, 0), edge.label, font=small_font)
            pad_x = 10
            pad_y = 5
            draw.rounded_rectangle(
                (
                    label_x - (label_box[2] - label_box[0]) / 2 - pad_x,
                    label_y - pad_y,
                    label_x + (label_box[2] - label_box[0]) / 2 + pad_x,
                    label_y + (label_box[3] - label_box[1]) + pad_y,
                ),
                radius=10,
                fill="#FFFDF8",
                outline="#E3D7C8",
            )
            draw.text(
                (label_x - (label_box[2] - label_box[0]) / 2, label_y),
                edge.label,
                font=small_font,
                fill=MUTED,
            )

    for node in NODES:
        draw.rounded_rectangle(
            (node.x, node.y, node.x + node.w, node.y + node.h),
            radius=22,
            fill=node.fill,
            outline=node.stroke,
            width=3,
        )
        draw_multiline_text(
            draw,
            node.x + 16,
            node.y + 14,
            node.w - 32,
            node.title,
            node.body,
            node_title_font,
            body_font,
            node.font_color,
        )

    image.save(PNG_PATH)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int = 18,
    color: RGBColor = RGBColor(35, 31, 28),
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for line in text.split("\n"):
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_card(slide, left: float, top: float, width: float, height: float, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(220, 209, 194)
    shape.line.width = Pt(1)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    accent = RGBColor(182, 111, 59)
    dark = RGBColor(28, 24, 22)
    ink = RGBColor(37, 31, 28)
    muted = RGBColor(106, 95, 85)
    panel = RGBColor(252, 247, 239)

    def new_slide(kicker: str, title: str, subtitle: str):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(246, 241, 232)
        add_textbox(slide, 0.55, 0.38, 4.8, 0.35, kicker, 12, accent, True)
        add_textbox(slide, 0.55, 0.72, 7.8, 0.55, title, 24, dark, True)
        add_textbox(slide, 0.58, 1.16, 11.6, 0.38, subtitle, 13, muted)
        return slide

    def add_bullet_block(slide, left: float, top: float, width: float, lines: list[str], size: int = 15):
        add_textbox(slide, left, top, width, 0.3 + 0.32 * len(lines), "\n".join(f"• {line}" for line in lines), size, muted)

    slide = new_slide(
        "INKPI GLOBAL PROJECT DECK",
        "InkPi 全局项目汇报",
        "围绕完整项目展开：终端形态、全链路架构、评分体系、云端同步、训练与测试，而不是只讲某一个功能点。",
    )
    add_card(slide, 0.72, 1.9, 4.2, 3.9, panel)
    add_textbox(
        slide,
        0.98,
        2.18,
        3.6,
        2.8,
        "项目定位\n"
        "InkPi 是一个面向书法评测场景的完整系统。\n\n"
        "它不仅包含本地单字识别与评分，还包含正式展示端、调试端、云端存储、微信小程序和训练链路。",
        16,
        ink,
    )
    add_card(slide, 5.2, 1.9, 7.3, 4.95, panel)
    slide.shapes.add_picture(str(PNG_PATH), Inches(5.42), Inches(2.12), width=Inches(6.88))

    slide = new_slide(
        "01 / PROJECT SCOPE",
        "项目目标与系统边界",
        "先把这个项目到底包含什么、解决什么问题、交付了哪些面讲清楚。",
    )
    cards = [
        ("设备侧", "PyQt 正式界面、本地拍照评测、本地 SQLite。"),
        ("调试侧", "WebUI 拍照 / 上传 / 历史详情 / 原始 debug。"),
        ("云端侧", "Flask Cloud API、CloudDatabase、设备同步接口。"),
        ("移动侧", "微信小程序登录、历史、详情展示。"),
        ("训练侧", "质量分档、ONNX 模型训练与导出。"),
        ("测试侧", "单元测试、接口回归、云同步集成测试。"),
    ]
    for idx, (title, body) in enumerate(cards):
        left = 0.72 + (idx % 3) * 4.08
        top = 1.72 + (idx // 3) * 2.2
        add_card(slide, left, top, 3.7, 1.8, panel)
        add_textbox(slide, left + 0.2, top + 0.16, 1.5, 0.3, title, 17, dark, True)
        add_textbox(slide, left + 0.2, top + 0.56, 3.15, 0.9, body, 13, muted)

    slide = new_slide(
        "02 / ARCHITECTURE",
        "整体架构与模块关系",
        "这一页讲全局结构：从输入入口，到本地评测，再到落库、同步、云端、小程序和训练产物。",
    )
    add_card(slide, 0.62, 1.65, 7.75, 5.35, panel)
    slide.shapes.add_picture(str(PNG_PATH), Inches(0.8), Inches(1.9), width=Inches(7.35))
    add_card(slide, 8.7, 1.65, 3.95, 5.35, panel)
    add_textbox(slide, 8.95, 1.95, 3.3, 0.3, "关键分层", 17, dark, True)
    add_bullet_block(
        slide,
        8.95,
        2.35,
        3.2,
        [
            "输入层：PyQt、WebUI、图像上传、小程序。",
            "运行层：Preprocessing、OCR、主分评分、四维解释分。",
            "结果层：EvaluationResult、本地库、Qt/WebUI 展示。",
            "同步层：CloudSync、Cloud API、CloudDatabase。",
            "训练层：quality_scorer 模型训练与导出。",
        ],
        14,
    )

    slide = new_slide(
        "03 / END-TO-END FLOW",
        "端到端业务流程",
        "这一页专门讲业务主链，不展开文件细节，只讲评测过程和数据怎么流动。",
    )
    blocks = [
        ("输入", "拍照或上传单字图像"),
        ("预处理", "预检、校正、主体提取"),
        ("识别", "LocalOcrService 输出 character + confidence"),
        ("主分评分", "QualityScorerService 输出 total_score 与 quality_level"),
        ("解释层", "DimensionScorerService 输出四维解释分"),
        ("结果落库", "EvaluationResult -> SQLite"),
        ("云端同步", "CloudSyncService -> Cloud API"),
        ("多端展示", "Qt / WebUI / MiniApp 读取结果"),
    ]
    for idx, (title, body) in enumerate(blocks):
        left = 0.75 + idx * 1.55
        add_card(slide, left, 2.25, 1.35, 2.5, panel)
        add_textbox(slide, left + 0.12, 2.48, 1.1, 0.3, title, 16, dark, True, PP_ALIGN.CENTER)
        add_textbox(slide, left + 0.12, 3.02, 1.1, 1.1, body, 12, muted, False, PP_ALIGN.CENTER)
        if idx < len(blocks) - 1:
            chevron = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + 1.36),
                Inches(3.15),
                Inches(0.16),
                Inches(0.35),
            )
            chevron.fill.solid()
            chevron.fill.fore_color.rgb = accent
            chevron.line.fill.background()
    add_textbox(
        slide,
        0.82,
        5.3,
        11.8,
        0.6,
        "现在项目的正式口径是单链路评测：不再依赖模板库主导评分，也不再把四维解释分当作第二套主分系统。",
        14,
        muted,
    )

    slide = new_slide(
        "04 / ASSESSMENT",
        "评测流程与主分生成机制",
        "这一页专门讲“系统怎么评测”，把输入、识别、特征、概率和总分校准逻辑讲清楚。",
    )
    add_card(slide, 0.72, 1.72, 5.85, 4.95, panel)
    add_textbox(slide, 0.95, 1.98, 2.8, 0.32, "评测流程", 17, dark, True)
    add_bullet_block(
        slide,
        0.95,
        2.38,
        5.05,
        [
            "PreprocessingService 先做预检、透视校正、红色网格去除、二值化和单字主体提取，确保送入模型的是稳定 ROI。",
            "同一次评测会同时产出 processed_image 和 ocr_image：前者服务主分模型，后者服务 OCR 识别。",
            "LocalOcrService 输出 character 与 ocr_confidence，既决定“识别出了什么字”，也给后续评分提供置信信号。",
            "QualityScorerService 读取 ONNX 输出，同时复用几何特征和图像质量特征，避免只靠类别概率硬判。",
            "主链最后统一装配成 total_score、quality_level、quality_confidence，并写入本地库、历史页和云同步 payload。",
        ],
        13,
    )
    add_card(slide, 6.82, 1.72, 5.78, 4.95, panel)
    add_textbox(slide, 7.05, 1.98, 3.4, 0.32, "主分 total_score 怎么来", 17, dark, True)
    add_bullet_block(
        slide,
        7.05,
        2.38,
        5.0,
        [
            "ONNX 先输出 bad / medium / good 三档概率，最大概率对应 quality_level，但类别本身不是最终分数。",
            "quality_features 包含 fg_ratio、bbox_ratio、center_quality、component_norm、edge_touch、texture_std，负责刻画墨量、居中、边缘触碰和纹理稳定性。",
            "calibration 中间量包含 feature_quality、probability_margin、quality_confidence_norm、score_range_fit，用来把“分类结果”变成“连续分数”。",
            "系统先定档，再在 bad 44-68 / medium 66-84 / good 82-98 区间内做连续校准，所以同一档内部也能拉开差距。",
            "如果模型直接给 raw_score 且 raw_score <= 1，则走 35% raw + 65% calibrated 的融合策略，兼顾端到端预测和规则校准。",
            "因此 total_score 的口径是“模型概率 + 图像特征 + 分档校准”的综合结果，而不是简单映射分。",
        ],
        13,
    )

    slide = new_slide(
        "05 / EVALUATION",
        "四维解释分与评价反馈生成",
        "这一页专门讲“系统怎么评价”，把四维解释分和最终反馈文案的生成口径讲清楚。",
    )
    add_card(slide, 0.72, 1.72, 5.85, 4.95, panel)
    add_textbox(slide, 0.95, 1.98, 3.0, 0.32, "四维解释分", 17, dark, True)
    add_bullet_block(
        slide,
        0.95,
        2.38,
        5.05,
        [
            "结构：重点看 center_quality、bbox_ratio_band、projection_balance、bbox_fill_band，回答“字是不是站得正、摆得稳”。",
            "笔画：重点看 texture_std_band、orientation_concentration、ink_ratio_band、component_norm_band，回答“线条控制是否均衡”。",
            "完整：重点看 ocr_confidence、dominant_share、(1-edge_touch)、subject_edge_safe，回答“有没有断裂、粘连、出框或缺损”。",
            "稳定：重点看 quality_confidence、probability_margin、feature_quality、score_range_fit，回答“这次判断是否可靠、是否落在合理区间”。",
            "四维分全部输出 0-100 的整数，只负责解释主分，不反推 total_score，也不替代 quality_level。",
        ],
        13,
    )
    add_card(slide, 6.82, 1.72, 5.78, 4.95, panel)
    add_textbox(slide, 7.05, 1.98, 3.2, 0.32, "评价与反馈生成", 17, dark, True)
    add_bullet_block(
        slide,
        7.05,
        2.38,
        5.0,
        [
            "EvaluationService 把 OCR、主分、四维解释层、调试中间量统一装配成 EvaluationResult，形成单次评测的标准结果对象。",
            "feedback 仍然保留单字符串，由 quality_level 对应模板池生成，再结合 total_score 选择更贴近当前水平的文案。",
            "如果 OCR 已识别出汉字，系统会把“识别字 + 当前评测等级”拼进最终反馈，让用户知道系统评的是哪一个字、处在什么水平。",
            "界面层会进一步提炼 strongest / weakest 维度，生成“最强项 / 待提升项”摘要，但不再发明第二套长文案评分系统。",
            "Qt 结果页强调结果展示，WebUI 详情页额外展开 score_debug 方便调参，MiniApp 只展示用户可理解的评价层。",
        ],
        13,
    )

    slide = new_slide(
        "06 / SURFACES",
        "前端与多终端角色分工",
        "同一个结果对象在不同端承担不同职责：正式展示、调试验证、云端查看、移动端查看。",
    )
    preview_specs = [
        ("qt_home_480x320.png", 0.72, "Qt 首页"),
        ("qt_result_480x320.png", 4.3, "Qt 结果页"),
        ("qt_history_480x320.png", 7.88, "Qt 历史页"),
    ]
    for name, left, label in preview_specs:
        path = SCREENSHOT_DIR / name
        add_card(slide, left, 1.75, 3.15, 4.55, panel)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(left + 0.12), Inches(1.92), width=Inches(2.92))
        add_textbox(slide, left + 0.12, 5.95, 2.9, 0.28, label, 13, muted, True, PP_ALIGN.CENTER)
    add_card(slide, 11.0, 1.75, 1.55, 4.55, panel)
    add_textbox(
        slide,
        11.12,
        2.0,
        1.2,
        3.7,
        "WebUI\n调试端\n\nCloud API\n云端接口\n\nMiniApp\n轻量查看",
        14,
        dark,
        True,
        PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        0.78,
        6.45,
        11.7,
        0.35,
        "Qt 负责正式展示，WebUI 负责 debug 细节，小程序负责轻量查看，云端负责统一结果入口和共享历史。",
        13,
        muted,
    )

    slide = new_slide(
        "07 / SCORING & TRAINING",
        "评分体系与训练链路",
        "评分是项目的一部分，而不是全部，但它是项目的核心能力之一，所以单独讲一页。",
    )
    add_card(slide, 0.72, 1.75, 5.8, 4.75, panel)
    add_textbox(slide, 0.95, 2.02, 3.2, 0.32, "运行时评分", 17, dark, True)
    add_bullet_block(
        slide,
        0.95,
        2.45,
        5.1,
        [
            "QualityScorerService 输出 total_score / quality_level / quality_confidence。",
            "保留 probabilities、quality_features、calibration 作为解释基础。",
            "DimensionScorerService 输出 structure / stroke / integrity / stability。",
            "主分 total_score 继续保留为官方分，四维分只负责解释和教学展示。",
        ],
        14,
    )
    add_card(slide, 6.78, 1.75, 5.8, 4.75, panel)
    add_textbox(slide, 7.02, 2.02, 3.2, 0.32, "训练链路", 17, dark, True)
    add_bullet_block(
        slide,
        7.02,
        2.45,
        5.1,
        [
            "build_quality_manifest.py 构建质量分档数据。",
            "train_quality_scorer.py 训练主分模型。",
            "导出 quality_scorer.onnx 与 metrics.json。",
            "训练产物直接回到运行时主链，服务端按当前模型输出做校准。",
        ],
        14,
    )

    slide = new_slide(
        "08 / DATA & CLOUD",
        "数据存储、同步与云端接口",
        "这一页讲项目的持久化能力和跨端一致性，而不是只讲单次评测。",
    )
    add_card(slide, 0.72, 1.75, 3.8, 4.9, panel)
    add_textbox(slide, 0.94, 2.0, 2.0, 0.3, "本地存储", 17, dark, True)
    add_bullet_block(
        slide,
        0.94,
        2.4,
        3.2,
        [
            "DatabaseService -> SQLite",
            "保存 total_score、feedback、character_name、quality_confidence",
            "新增 dimension_scores_json / score_debug_json",
            "老记录保持 null，不强制回填",
        ],
        14,
    )
    add_card(slide, 4.8, 1.75, 3.8, 4.9, panel)
    add_textbox(slide, 5.02, 2.0, 2.0, 0.3, "同步与云端", 17, dark, True)
    add_bullet_block(
        slide,
        5.02,
        2.4,
        3.2,
        [
            "CloudSyncService 异步上传结果",
            "Cloud API 统一处理设备上报",
            "CloudDatabase 保存共享历史",
            "列表接口轻量返回，详情接口返回 debug",
        ],
        14,
    )
    add_card(slide, 8.88, 1.75, 3.7, 4.9, panel)
    add_textbox(slide, 9.1, 2.0, 2.0, 0.3, "小程序 / 查询", 17, dark, True)
    add_bullet_block(
        slide,
        9.1,
        2.4,
        3.1,
        [
            "MiniApp 读取云端历史与详情",
            "历史列表保持轻量",
            "详情页展示主分、反馈和四维解释分",
            "不展示 score_debug 原始数据",
        ],
        14,
    )

    slide = new_slide(
        "09 / ENGINEERING STATUS",
        "工程化、测试与部署状态",
        "最后一页回到项目交付：它现在到底到什么程度，哪些地方已经稳定，哪些地方后面继续做。",
    )
    add_card(slide, 0.72, 1.72, 5.9, 4.95, panel)
    add_textbox(slide, 0.95, 1.98, 2.4, 0.3, "当前已完成", 17, dark, True)
    add_bullet_block(
        slide,
        0.95,
        2.38,
        5.2,
        [
            "本地单链路评测已跑通。",
            "Qt / WebUI / Cloud API / MiniApp 的结果结构已统一。",
            "draw.io 总流程图与项目 PPT 已和当前代码状态对齐。",
            "单元测试、接口测试、云同步集成测试已通过。",
        ],
        14,
    )
    add_card(slide, 6.88, 1.72, 5.72, 4.95, panel)
    add_textbox(slide, 7.12, 1.98, 2.4, 0.3, "下一步优化", 17, dark, True)
    add_bullet_block(
        slide,
        7.12,
        2.38,
        5.0,
        [
            "继续优化书法审美相关特征和四维解释口径。",
            "补真实移动端和 WebUI 的人工演示截图。",
            "按树莓派真机环境继续补联调与部署说明。",
            "把答辩版进一步压缩成更短的展示节奏。",
        ],
        14,
    )

    prs.save(PPTX_PATH)


def main() -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    build_drawio()
    build_png()
    build_ppt()
    print(DRAWIO_PATH)
    print(PNG_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
