from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


DESKTOP = Path(r"C:\Users\zongrui\Desktop")
OUTPUT_PATH = DESKTOP / "InkPi_比赛答辩版.pptx"
PHOTO_DIR = DESKTOP / "photos"

PRIMARY = RGBColor(13, 17, 29)
PANEL = RGBColor(24, 30, 49)
PANEL_SOFT = RGBColor(29, 36, 59)
TEXT = RGBColor(244, 246, 252)
MUTED = RGBColor(160, 171, 196)
ACCENT = RGBColor(233, 179, 88)
ACCENT_SOFT = RGBColor(73, 197, 241)
GREEN = RGBColor(76, 187, 125)
RED = RGBColor(225, 104, 99)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = 13.333
SLIDE_H = 7.5

SERVER_RESULTS = [
    ("黄", 93),
    ("京", 96),
    ("濟", 93),
    ("水", 94),
    ("序", 94),
    ("洛", 91),
    ("三", 96),
    ("年", 95),
    ("并", 96),
    ("賦", 94),
    ("师", 95),
    ("還", 95),
    ("朝", 94),
    ("初", 96),
    ("余", 91),
]


def set_bg(slide, color=PRIMARY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_blob(slide, left, top, width, height, color, transparency=0.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    fill.transparency = transparency
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title=None, fill_color=PANEL, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.adjustments[0] = 0.12
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(1.2)
    else:
        line.fill.background()

    if title:
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.22),
            Inches(top + 0.18),
            Inches(width - 0.44),
            Inches(0.3),
        )
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, lines, color=TEXT, size=18, bullet_color=ACCENT_SOFT):
    step = 0.42
    for idx, line in enumerate(lines):
        cy = top + idx * step
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(cy + 0.06),
            Inches(0.12),
            Inches(0.12),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color
        dot.line.fill.background()
        add_text(slide, left + 0.18, cy, width - 0.18, 0.26, line, size=size, color=color)


def add_metric_card(slide, left, top, width, height, value, label, note, accent):
    add_card(slide, left, top, width, height, fill_color=PANEL_SOFT)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, left + 0.22, top + 0.18, width - 0.44, 0.48, value, size=28, color=WHITE, bold=True)
    add_text(slide, left + 0.22, top + 0.72, width - 0.44, 0.28, label, size=13, color=MUTED, bold=True)
    add_text(slide, left + 0.22, top + 1.02, width - 0.44, 0.48, note, size=12, color=TEXT)


def add_section_title(slide, kicker, title, subtitle):
    add_text(slide, 0.7, 0.52, 3.0, 0.2, kicker.upper(), size=11, color=ACCENT, bold=True)
    add_text(slide, 0.7, 0.78, 8.8, 0.75, title, size=28, color=WHITE, bold=True)
    add_text(slide, 0.72, 1.45, 8.8, 0.45, subtitle, size=13, color=MUTED)


def add_footer(slide, text="InkPi | 自动书法评测 Demo"):
    add_text(slide, 0.7, 7.0, 4.5, 0.2, text, size=10, color=MUTED)
    add_text(slide, 11.6, 7.0, 1.0, 0.2, "答辩版", size=10, color=ACCENT, align=PP_ALIGN.RIGHT, bold=True)


def add_flow_node(slide, left, top, width, title, body, accent_color):
    add_card(slide, left, top, width, 1.32, fill_color=PANEL)
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.16),
        Inches(top + 0.14),
        Inches(0.72),
        Inches(0.24),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = accent_color
    tag.line.fill.background()
    add_text(slide, left + 0.19, top + 0.11, 0.66, 0.2, title, size=11, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.16, top + 0.48, width - 0.32, 0.62, body, size=12, color=TEXT)


def add_arrow(slide, left, top, width=0.45):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.36),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background()


def add_score_row(slide, left, top, width, char, score):
    add_text(slide, left, top, 0.38, 0.2, char, size=17, color=WHITE, bold=True)
    rail = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.42),
        Inches(top + 0.03),
        Inches(width - 0.98),
        Inches(0.17),
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = RGBColor(54, 61, 83)
    rail.line.fill.background()
    fill_width = (width - 0.98) * max(0.12, min(1.0, (score - 80) / 20))
    fill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.42),
        Inches(top + 0.03),
        Inches(fill_width),
        Inches(0.17),
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = ACCENT if score >= 95 else ACCENT_SOFT if score >= 93 else GREEN
    fill.line.fill.background()
    add_text(slide, left + width - 0.42, top - 0.02, 0.4, 0.2, str(score), size=14, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)


def maybe_add_image(slide, path: Path, left, top, width, height):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    sample_1 = PHOTO_DIR / "403ec733-b201-422d-8403-e1a1d0879a6b.png"
    sample_2 = PHOTO_DIR / "479aeeb7-8a79-4c8d-9aab-01ebae826887.png"

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_blob(slide, 9.5, -0.8, 4.0, 4.0, ACCENT, 0.78)
    add_blob(slide, -0.9, 5.5, 3.8, 3.8, ACCENT_SOFT, 0.87)
    add_text(slide, 0.72, 0.58, 2.2, 0.2, "INKPI | COMPETITION PITCH", size=11, color=ACCENT, bold=True)
    add_text(slide, 0.7, 1.0, 6.4, 0.9, "InkPi\n书法自动评测 Demo", size=28, color=WHITE, bold=True)
    add_text(
        slide,
        0.74,
        2.1,
        5.9,
        0.7,
        "Figma 风格答辩版表达：深色舞台、关键指标先行、\n把“能跑”“能识别”“能同步”讲成一条完整产品链。",
        size=14,
        color=MUTED,
    )
    add_metric_card(slide, 0.72, 3.0, 2.0, 1.55, "15 / 15", "图片回归通过", "服务器真机接口批量验证", ACCENT)
    add_metric_card(slide, 2.92, 3.0, 2.0, 1.55, "1 条", "自动评测主链", "预处理 -> OCR -> ONNX 评分", ACCENT_SOFT)
    add_metric_card(slide, 5.12, 3.0, 2.0, 1.55, "WebUI", "当前演示形态", "后端服务器已可直接展示", GREEN)
    add_card(slide, 8.2, 0.92, 4.4, 5.6, fill_color=PANEL, line_color=RGBColor(58, 65, 94))
    maybe_add_image(slide, sample_1, 8.45, 1.16, 1.95, 4.9)
    maybe_add_image(slide, sample_2, 10.42, 1.16, 1.95, 4.9)
    add_text(slide, 8.45, 6.16, 3.9, 0.32, "真实书法样张 | 当前服务器 OCR + ONNX 评测实测通过", size=11, color=MUTED)
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "01 / POSITIONING", "项目定位与比赛价值", "把树莓派设备端、云端历史与微信小程序串成一个可讲、可演示、可落地的完整作品。")
    add_card(slide, 0.72, 2.0, 5.8, 4.35, title="为什么这个方向值得做")
    add_bullets(
        slide,
        0.98,
        2.55,
        5.1,
        [
            "传统书法展示偏静态，评委很难一眼看到“系统能力”与“产品闭环”。",
            "InkPi 把拍照、识别、评分、存档、同步全部压缩到一次交互里。",
            "硬件端可做展台 Demo，云端和小程序又能证明项目具备延展性。",
            "比赛讲法不再是“算法实验”，而是“可以落地的交互产品”。",
        ],
    )
    add_card(slide, 6.78, 2.0, 5.82, 4.35, title="当前成品形态")
    add_bullets(
        slide,
        7.04,
        2.55,
        5.05,
        [
            "设备端：本地拍照、自动识别、自动评分。",
            "服务器端：当前承载 WebUI 展示、云端接口与历史结果。",
            "小程序端：登录查看历史成绩，补充移动端展示入口。",
            "比赛现场：优先展示“单链路自动评测”而不是复杂配置。 ",
        ],
        bullet_color=ACCENT,
    )
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "02 / ARCHITECTURE", "系统架构", "把设备端、后端服务、小程序与云端历史同步成一张图，强调当前可运行的主链路。")
    add_card(slide, 0.72, 2.0, 3.6, 4.5, title="设备 / WebUI")
    add_bullets(slide, 0.98, 2.55, 3.0, ["拍照或上传书法图片", "预处理提取单字主体", "本地 PaddleOCR 自动识别", "ONNX 评分模型返回总分与等级"])
    add_card(slide, 4.86, 2.0, 3.62, 4.5, title="后端 / 数据")
    add_bullets(slide, 5.12, 2.55, 3.0, ["Flask WebUI 提供评测接口", "SQLite 保存本地历史结果", "云端 API 聚合设备上传记录", "支持后续模型迭代与版本切换"])
    add_card(slide, 9.0, 2.0, 3.6, 4.5, title="展示 / 小程序")
    add_bullets(slide, 9.26, 2.55, 3.0, ["微信小程序登录后查看历史", "树莓派当前不可用时，服务器可直接展示前端", "适合比赛答辩与现场观众同步查看结果", "便于继续扩展设备联网能力"])
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "03 / PIPELINE", "自动评测主链路", "这一版已经去掉模板锁字和双模式绕路，评测流程只有一条线。")
    add_flow_node(slide, 0.74, 2.36, 2.55, "步骤 1", "图像预处理\n去红格、提主体、保留适合 OCR 的轻裁剪灰度图", ACCENT)
    add_arrow(slide, 3.42, 2.82)
    add_flow_node(slide, 3.92, 2.36, 2.55, "步骤 2", "本地官方 OCR\nPaddleOCR 直接识别汉字，不再手动锁字", ACCENT_SOFT)
    add_arrow(slide, 6.60, 2.82)
    add_flow_node(slide, 7.10, 2.36, 2.55, "步骤 3", "ONNX 评分模型\n输出质量等级，再由服务端校准成连续分数", GREEN)
    add_arrow(slide, 9.78, 2.82)
    add_flow_node(slide, 10.28, 2.36, 2.3, "步骤 4", "结果展示与同步\n本地显示、数据库落盘、云端接口上传", ACCENT)
    add_card(slide, 0.74, 4.45, 11.84, 1.3, fill_color=PANEL_SOFT, line_color=RGBColor(70, 80, 110))
    add_text(slide, 1.02, 4.8, 11.0, 0.32, "当前主链已经明确删除：模板库主导评分、Siamese 对图比较、手动锁定评测字、综合/兜底双模式。", size=15, color=WHITE, bold=True)
    add_text(slide, 1.02, 5.22, 11.0, 0.24, "评测表达更清晰：识别字 + 置信度 + 总分 + 等级 + 简短反馈。", size=12, color=MUTED)
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "04 / EVIDENCE", "评价与评分依据", "把“为什么这么打分”讲成可解释的工程逻辑，而不是口头形容词。")
    add_card(slide, 0.72, 2.0, 4.6, 4.55, title="识别依据")
    add_bullets(
        slide,
        0.98,
        2.55,
        4.0,
        [
            "单字主体先经过专门的 OCR 轻裁剪支路。",
            "官方中文 OCR 在本地设备运行，低置信度时直接拒识。",
            "Web 常驻服务已补上串行锁与异常后自动重建，引擎稳定性更高。",
        ],
        bullet_color=ACCENT_SOFT,
    )
    add_card(slide, 5.56, 2.0, 7.04, 4.55, title="评分依据")
    add_text(slide, 5.84, 2.56, 6.4, 0.24, "当前总分 = 等级区间 + 图像特征校准", size=18, color=WHITE, bold=True)
    add_bullets(
        slide,
        5.84,
        3.02,
        6.2,
        [
            "特征分权重：前景占比 28%、重心质量 24%、连通域复杂度 16%、触边率 16%、纹理波动 16%。",
            "置信分权重：模型最大类别概率 60%、第一第二名间隔 25%、OCR 置信度 15%。",
            "最终融合：特征分 72%、置信分 28%。",
            "等级区间：bad 44-68、medium 66-84、good 82-98。",
        ],
        bullet_color=ACCENT,
    )
    add_text(slide, 5.84, 5.1, 6.2, 0.45, "这次已经从“只要判成 good 就几乎固定 92 分”改成“先定分数段，再按图像质量做段内连续分”。", size=12, color=MUTED)
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "05 / RESULTS", "当前服务器实测结果", "同一批真实书法样图已经完成 15 / 15 自动识别评测，分数也从单一 92 分拉开。")
    add_card(slide, 0.72, 2.0, 12.0, 4.62, title="Desktop/photos 实测分布")
    add_text(slide, 0.98, 2.48, 4.8, 0.25, "服务器 WebUI 真实接口批量测试", size=14, color=WHITE, bold=True)
    add_text(slide, 0.98, 2.82, 4.8, 0.25, "结果：15 / 15 通过，分数区间 91 - 96", size=12, color=MUTED)
    left_x = 1.0
    right_x = 6.9
    y0 = 3.25
    for idx, (char, score) in enumerate(SERVER_RESULTS[:8]):
        add_score_row(slide, left_x, y0 + idx * 0.38, 4.8, char, score)
    for idx, (char, score) in enumerate(SERVER_RESULTS[8:]):
        add_score_row(slide, right_x, y0 + idx * 0.38, 4.8, char, score)
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "06 / DEMO FLOW", "比赛现场怎么讲", "把答辩节奏压成 4 句话：采集、识别、评分、同步，讲产品闭环而不是讲脚本细节。")
    add_card(slide, 0.72, 2.0, 12.0, 4.5, title="建议演示台词")
    add_bullets(
        slide,
        1.0,
        2.55,
        10.8,
        [
            "第一句：这是一个面向书法展示场景的自动评测系统，当前服务已经可以直接跑在服务器前端上。",
            "第二句：上传书法图片后，系统先做预处理，再用本地 OCR 自动识别汉字。",
            "第三句：识别结果进入 ONNX 评分模型，系统返回总分、等级和反馈，而不是固定模板匹配结果。",
            "第四句：评测结果会落到本地数据库，并同步到云端和微信小程序，形成历史可追踪能力。",
        ],
        bullet_color=ACCENT_SOFT,
    )
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_section_title(slide, "07 / NEXT", "下一阶段优化方向", "这份答辩版先把“能演示”讲清楚；后面继续把评分模型本体做得更像书法审美判断。")
    add_card(slide, 0.72, 2.0, 5.9, 4.4, title="已经完成")
    add_bullets(
        slide,
        0.98,
        2.55,
        5.1,
        [
            "服务器端 WebUI 已可直接展示。",
            "15 张真实样图 OCR 识别 15 / 15 通过。",
            "评分已从固定 92 分改为连续分布。",
            "小程序、云端同步和历史记录链路都已具备。 ",
        ],
    )
    add_card(slide, 6.82, 2.0, 5.78, 4.4, title="继续优化")
    add_bullets(
        slide,
        7.08,
        2.55,
        5.0,
        [
            "重训评分模型本体，进一步拉开真实名家字之间的分数差距。",
            "补充更强的书法风格标签与评价解释，提升答辩说服力。",
            "树莓派恢复后，把同一套前端和模型重新回灌到设备端实机。",
            "把小程序视觉也同步升级成统一比赛展示风格。",
        ],
        bullet_color=ACCENT,
    )
    add_footer(slide, "InkPi | 比赛答辩版 Deck")

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    out = build_presentation()
    print(out)
